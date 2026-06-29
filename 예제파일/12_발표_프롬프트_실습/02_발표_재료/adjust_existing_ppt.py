import os
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR

# UTF-8 출력 재설정 (윈도우 한글 인코딩 깨짐 방지)
sys.stdout.reconfigure(encoding='utf-8')

# 색상 상수 정의 (가이드 준수 및 고품질 디자인 적용)
COLOR_BG = RGBColor(255, 255, 255)         # 배경 흰색
COLOR_PRIMARY = RGBColor(43, 87, 151)       # 파랑 (#2B5797)
COLOR_SECONDARY = RGBColor(232, 119, 34)    # 포인트 주황 (#E87722)
COLOR_DARK = RGBColor(43, 87, 151)          # 메인 글씨 색상 (파랑 통일)
COLOR_MUTED = RGBColor(120, 130, 145)       # 보조 그레이 색상 (푸터/쪽번호용)
COLOR_CARD_BG = RGBColor(248, 249, 250)     # 카드 배경용 연한 회색
COLOR_WHITE = RGBColor(255, 255, 255)

def clean_slide(slide):
    """슬라이드의 모든 기존 요소를 삭제하고 배경을 흰색으로 초기화"""
    # 1. 배경 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG
    
    # 2. 모든 Shape 삭제 (안전하게 루프 돌기 위해 리스트로 복사 후 제거)
    shapes = list(slide.shapes)
    for s in shapes:
        sp = s._element
        sp.getparent().remove(sp)

def add_header(slide, title_text, subtitle_text):
    """헤더 추가 (2cm 이상 상단 여백 고려: 0.8 inches = 2.03cm)"""
    # 좌측 파란색 데코 라인
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(0.12), Inches(0.8))
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR_PRIMARY
    deco.line.fill.background()
    
    # 타이틀 박스
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.42), Inches(11.5), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = '맑은 고딕'
    p.font.size = Pt(28)  # 28pt 이상, 굵게
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 서브 메시지 박스 (주황색 포인트 활용)
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.92), Inches(11.5), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_top = tf_sub.margin_bottom = tf_sub.margin_left = tf_sub.margin_right = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.name = '맑은 고딕'
    p_sub.font.size = Pt(15)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_SECONDARY
    
    # 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LINE = RGBColor(220, 225, 230)
    line.line.fill.background()

def add_footer_and_pagenum(slide, page_num):
    """슬라이드 하단에 부서명("영업기획팀") 및 쪽번호 추가 (2cm 여백 고려: 6.9 inches)"""
    # 부서명 (좌측)
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(4.0), Inches(0.3))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "영업기획팀"
    p_l.font.name = '맑은 고딕'
    p_l.font.size = Pt(11)
    p_l.font.color.rgb = COLOR_MUTED
    
    # 쪽번호 (우측)
    right_box = slide.shapes.add_textbox(Inches(11.0), Inches(6.9), Inches(1.5), Inches(0.3))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = str(page_num)
    p_r.font.name = '맑은 고딕'
    p_r.font.size = Pt(11)
    p_r.font.color.rgb = COLOR_MUTED
    p_r.alignment = PP_ALIGN.RIGHT

def format_text_box(slide, text_lines, left, top, width, height, font_size=18, bg_color=None):
    """3줄 이하 텍스트 본문 생성 상자"""
    card = None
    if bg_color:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = RGBColor(230, 235, 240)
        card.line.width = Pt(1)
        
    t_box = slide.shapes.add_textbox(left + Inches(0.15) if bg_color else left, 
                                     top + Inches(0.15) if bg_color else top, 
                                     width - Inches(0.3) if bg_color else width, 
                                     height - Inches(0.3) if bg_color else height)
    tf = t_box.text_frame
    tf.word_wrap = True
    
    for idx, line in enumerate(text_lines):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = line
        p.font.name = '맑은 고딕'
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_DARK
        p.font.bold = False
        p.space_after = Pt(12)
        
    return card

def add_notes(slide, notes_text):
    """슬라이드 노트 등록"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def main():
    pptx_path = "c:/Users/user/Desktop/ai사무자동화/예제파일/12_발표_프롬프트_실습/01_보고서_원본/상반기_매출_발표자료.pptx"
    
    prs = Presentation(pptx_path)
    # 16:9 와이드스크린으로 일괄 강제 조정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # SLIDE 1: 표지
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    clean_slide(slide1)
    
    # 세련된 표지 디자인 (흰색 배경 + 파란색 포인트 카드 형태)
    cover_card = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.2), Inches(4.5))
    cover_card.fill.solid()
    cover_card.fill.fore_color.rgb = COLOR_PRIMARY
    cover_card.line.fill.background()
    
    title_box = slide1.shapes.add_textbox(Inches(1.3), Inches(2.2), Inches(11.0), Inches(1.5))
    tf1 = title_box.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "2025년 상반기 매출 보고"
    p1.font.name = '맑은 고딕'
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    
    sub_box = slide1.shapes.add_textbox(Inches(1.3), Inches(3.5), Inches(11.0), Inches(1.0))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "보고기간: 2025년 1월 ~ 6월  |  부서: 영업기획팀"
    p_sub.font.name = '맑은 고딕'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_SECONDARY
    
    info_box = slide1.shapes.add_textbox(Inches(1.3), Inches(4.8), Inches(11.0), Inches(1.0))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "작성자: 김민서 대리  |  작성일: 2025년 7월 5일"
    p_info.font.name = '맑은 고딕'
    p_info.font.size = Pt(14)
    p_info.font.color.rgb = COLOR_MUTED
    
    add_notes(slide1, 
        "안녕하십니까. 영업기획팀 김민서 대리입니다. 지금부터 2025년 상반기 매출 보고 발표를 시작하겠습니다. "
        "본 보고는 지난 1월부터 6월까지의 당사 영업 실적을 체계적으로 점검하여 핵심 발견사항을 제시하기 위해 마련되었습니다. "
        "다음 슬라이드로 이동해 전체 발표 목차를 안내해 드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 2: 목차 (3~5개 유지 가이드 반영: 5개 항목 병합 구성)
    # -------------------------------------------------------------
    slide2 = prs.slides[1]
    clean_slide(slide2)
    add_header(slide2, "목차", "상반기 실적 분석부터 전략 제안까지 5단계 흐름으로 보고드리겠습니다.")
    add_footer_and_pagenum(slide2, 2)
    
    # 5개 카드 가로 및 세로 배치
    toc_items = [
        ("01", "전체 실적 요약"),
        ("02", "월별 및 제품별 분석"),
        ("03", "지역별 매출 분포"),
        ("04", "주요 성과 및 이슈"),
        ("05", "하반기 전략 제안")
    ]
    
    col_w = Inches(11.7)
    row_h = Inches(0.7)
    start_x = Inches(0.8)
    start_y = Inches(1.7)
    
    for idx, (num, title) in enumerate(toc_items):
        ty = start_y + idx * (row_h + Inches(0.2))
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, ty, col_w, row_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = RGBColor(230, 235, 240)
        card.line.width = Pt(1)
        
        # 텍스트
        t_box = slide2.shapes.add_textbox(start_x + Inches(0.2), ty + Inches(0.12), col_w - Inches(0.4), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num}.  {title}"
        p.font.name = '맑은 고딕'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
    add_notes(slide2, 
        "이 슬라이드에서 말씀드릴 것은 오늘 발표할 5가지 핵심 목차입니다. "
        "구체적으로 말씀드리면, 전체 실적 요약 보고를 시작으로 월별/제품별 상세 데이터, 지역별 분포를 종합 분석하겠습니다. "
        "이어서 주요 성과와 과제를 진단하고, 최종적으로 하반기 4대 영업 전략을 제안하는 순으로 마치겠습니다. "
        "먼저 전체 실적 요약을 보고드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 3: 전체 실적 요약 (글 3줄 이내 준수)
    # -------------------------------------------------------------
    slide3 = prs.slides[2]
    clean_slide(slide3)
    add_header(slide3, "1. 전체 실적 요약", "상반기 총매출 28억 4,200만 원 달성, 전년 동기 대비 +12.3% 큰 폭 성장")
    add_footer_and_pagenum(slide3, 3)
    
    # 3줄 텍스트 본문 생성
    summary_text = [
        "· 실적 지표: 상반기 총매출 28억 4,200만 원 (전년 동기 대비 +12.3% 증가)",
        "· 판매 효율: 총 판매 건수 1,847건 (전년비 +8.7%), 월평균 매출 4억 7,367만 원",
        "· 월별 극단: 4월 봄 프로모션 최고 실적(6.9억) / 3월 계절 비수기 최저 실적(2.08억)"
    ]
    format_text_box(slide3, summary_text, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), font_size=18, bg_color=COLOR_CARD_BG)
    
    add_notes(slide3, 
        "이 슬라이드에서 말씀드릴 것은 2025년 상반기 전체 영업 실적 요약입니다. "
        "구체적으로 말씀드리면, 상반기 총매출은 28억 4,200만 원으로 전년 동기 대비 12.3% 크게 증가하였으며, 총 판매 건수는 1,847건으로 8.7% 늘어났습니다. "
        "월별로는 봄 프로모션 효과를 얻은 4월이 6억 9,003만 원으로 최고치를 기록한 반면 비수기인 3월은 2억 790만 원으로 최저치에 머물렀습니다. "
        "다음 슬라이드로 이동하여 월별 매출 흐름을 꺾은선 차트로 더욱 상세히 확인해 보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 4: 월별 매출 현황 (꺾은선 차트 삽입)
    # -------------------------------------------------------------
    slide4 = prs.slides[3]
    clean_slide(slide4)
    add_header(slide4, "2. 월별 매출 현황", "시간 흐름에 따른 상반기 월 매출 추이 (3월 최저점 극복 후 4월 최고점 반등)")
    add_footer_and_pagenum(slide4, 4)
    
    # 차트 데이터 준비 (월별 합산 억 원 단위)
    chart_data = CategoryChartData()
    chart_data.categories = ['1월', '2월', '3월', '4월', '5월', '6월']
    # 매출액 (억 원): 1월: 4.74억, 2월: 7.21억, 3월: 2.08억, 4월: 6.90억, 5월: 4.90억, 6월: 3.07억
    chart_data.add_series('총매출액 (억 원)', (4.74, 7.21, 2.08, 6.90, 4.90, 3.07))
    
    # 꺾은선 차트 (Line Chart) 네이티브 삽입
    x, y, cx, cy = Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5)
    chart = slide4.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = False
    
    add_notes(slide4, 
        "이 슬라이드에서 말씀드릴 것은 월별 매출 변동 추이입니다. "
        "구체적으로 말씀드리면, 꺾은선 차트에서 나타나듯 1월 일시 하락 후 2월에 7억 2천만 원대까지 급증했습니다. "
        "이후 3월에 계절 비수기로 최저점을 찍었으나, 4월 봄 프로모션이 폭발적인 성과를 거두어 6억 9천만 원대로 재반등에 성공했습니다. "
        "다음 슬라이드에서는 이러한 매출을 이끈 제품별 비중 분석 결과를 보여드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 5: 제품별 분석 (원그래프 삽입)
    # -------------------------------------------------------------
    slide5 = prs.slides[4]
    clean_slide(slide5)
    add_header(slide5, "3. 제품별 분석", "프리미엄(44.2%)과 엔터프라이즈(39.1%) 제품군이 성장의 핵심 축 담당")
    add_footer_and_pagenum(slide5, 5)
    
    # 원그래프 데이터 준비
    chart_data_pie = CategoryChartData()
    chart_data_pie.categories = ['프리미엄', '엔터프라이즈', '베이직']
    chart_data_pie.add_series('매출 비중 (%)', (44.2, 39.1, 16.7))
    
    # 원그래프 (Pie Chart) 삽입
    x, y, cx, cy = Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5)
    chart_pie = slide5.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data_pie
    ).chart
    chart_pie.has_legend = True
    chart_pie.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    add_notes(slide5, 
        "이 슬라이드에서 말씀드릴 것은 제품군별 매출 비중 분석입니다. "
        "구체적으로 말씀드리면, 전체 매출 비중의 44.2%를 차지한 프리미엄 제품군이 매출 기여도 1위를 기록했습니다. "
        "엔터프라이즈는 B2B 영업 활성화로 39.1%의 매출 비중을 차지하여 2위를 기록했습니다. "
        "두 주력 제품군이 전체 매출의 83.3%를 주도하고 있음을 볼 수 있습니다. "
        "다음 슬라이드에서는 지역별 매출 점유율을 확인해 보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 6: 지역별 분석 (가로 막대 차트 삽입)
    # -------------------------------------------------------------
    slide6 = prs.slides[5]
    clean_slide(slide6)
    add_header(slide6, "4. 지역별 분석", "수도권(62.3%) 편중 현상 지속으로 전국 단위 커버리지 확장 과제 도출")
    add_footer_and_pagenum(slide6, 6)
    
    # 가로 막대 차트 데이터 준비
    chart_data_bar = CategoryChartData()
    # categories는 가로막대의 경우 아래부터 위 순서대로 표시되므로 하위 지역부터 정렬
    chart_data_bar.categories = ['기타 (3.4%)', '호남권 (6.2%)', '충청권 (9.4%)', '영남권 (18.7%)', '수도권 (62.3%)']
    chart_data_bar.add_series('매출 점유율 (%)', (3.4, 6.2, 9.4, 18.7, 62.3))
    
    # 가로 막대 차트 (Bar Clustered) 삽입
    x, y, cx, cy = Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5)
    chart_bar = slide6.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data_bar
    ).chart
    chart_bar.has_legend = False
    
    add_notes(slide6, 
        "이 슬라이드에서 말씀드릴 것은 지역별 상반기 매출 분포 순위입니다. "
        "구체적으로 말씀드리면, 가로 막대 차트에서 가장 긴 비중을 차지하는 수도권이 62.3%의 높은 점유율로 과반 이상을 지배하고 있습니다. "
        "뒤이어 영남권이 18.7%를 점유했고, 호남과 기타 지방권역은 각 6% 미만의 매우 저조한 침투율을 보였습니다. "
        "따라서 수도권 편중 해결을 위해 하반기 지방권 확장 영업이 요구됩니다. "
        "다음 슬라이드에서는 상반기 주요 성과 및 이슈를 종합 진단해 보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 7: 주요 성과 및 이슈 (성과/이슈 각 3줄 이내 레이아웃)
    # -------------------------------------------------------------
    slide7 = prs.slides[6]
    clean_slide(slide7)
    add_header(slide7, "5. 주요 성과 및 이슈", "유치 목표 초과 및 만족도 NPS 향상 vs 비수기 매출 저하 해결 시급")
    add_footer_and_pagenum(slide7, 7)
    
    # 좌우 2열 카드 레이아웃
    col_w = Inches(5.6)
    card_h = Inches(4.5)
    
    # 성과 카드 (좌측)
    format_text_box(slide7, [
        "■ 주요 성과 (Key Achievements)",
        "· 신규 고객 유치 목표 달성률 108% (302건 달성)",
        "· 고객 만족도 NPS 72점 달성 (전년비 7점 상승)",
        "· 엔터프라이즈 신규 3개사 수주 (연 5억+ 매출 효과)"
    ], Inches(0.8), Inches(1.8), col_w, card_h, font_size=14, bg_color=COLOR_CARD_BG)
    
    # 이슈 카드 (우측 - 주황 포인트 테두리)
    format_text_box(slide7, [
        "■ 극복 과제 (Key Issues)",
        "· 3월 분기 전환기 영업 공백에 따른 비수기 매출 급감 해결",
        "· 베이직 요금 고객 단가 하락 방지 부가서비스 패키징",
        "· 호남 및 강원 지역 침투율 저조 극복을 위한 파트너망 구축"
    ], Inches(6.9), Inches(1.8), col_w, card_h, font_size=14, bg_color=COLOR_CARD_BG)
    
    add_notes(slide7, 
        "이 슬라이드에서 말씀드릴 것은 상반기 성과 및 직면한 극복 과제입니다. "
        "구체적으로 말씀드리면, 신규 영업 목표 108% 달성과 함께 고객 충성도 NPS 지수가 72점까지 상승하는 뚜렷한 성과를 거두었습니다. "
        "하지만 3월 영업 공백 관리 체계 부족, 베이직 단가 하락세, 그리고 강원/호남 파트너 망 저조 등의 고질적인 이슈는 속히 극복해야 합니다. "
        "이 분석을 집대성하여 하반기 최종 제안 내용을 보고드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 8: 하반기 전략 (요약/제안 1장 형식 개편)
    # 가이드: 핵심 발견 3줄 / 제안, 다음 단계 2줄
    # -------------------------------------------------------------
    slide8 = prs.slides[7]
    clean_slide(slide8)
    add_header(slide8, "요약 및 제안", "성장세를 가속화하기 위한 핵심 분석 결과 및 하반기 전략 액션 제안")
    add_footer_and_pagenum(slide8, 8)
    
    # 핵심 발견 3줄 박스
    format_text_box(slide8, [
        "[핵심 발견 (Key Findings)]",
        "1. 상반기 총매출 28.4억 원(+12.3% YoY) 달성으로 전반적인 비즈니스 성장세 확인",
        "2. 프리미엄(44.2%)과 엔터프라이즈(39.1%) 제품군이 전체 성장의 중추적 역할 수행",
        "3. 수도권 편중(62.3%) 문제 및 3월 계절 비수기 매출 하락 제어를 위한 방안 마련 필요"
    ], Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2), font_size=13, bg_color=COLOR_CARD_BG)
    
    # 제안/다음 단계 2줄 박스
    format_text_box(slide8, [
        "[제안 및 다음 단계 (Proposals & Next Steps)]",
        "1. 프리미엄 프로모션 유치 및 엔터프라이즈 분기별 세미나를 통한 기존 고객 리텐션 확보",
        "2. 호남·강원 지역 파트너 계약 및 입찰 활성화, 솔루션 마케팅 예산 20% 증액 편성"
    ], Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.0), font_size=13, bg_color=COLOR_CARD_BG)
    
    add_notes(slide8, 
        "이 슬라이드에서 말씀드릴 것은 상반기 전체 요약 및 하반기 추진 제안입니다. "
        "구체적으로 말씀드리면, 매출 성장을 증명함과 동시에 주력 제품군 편중 극복과 지방 시장 활로 개척이라는 전략 과제를 발굴했습니다. "
        "이에 대한 제안으로 프리미엄 프로모션 유치, 지방 파트너쉽 강화, 디지털 광고 예산 20% 증액 조치를 적극 시행할 것을 건의드립니다. "
        "마지막으로 질의응답을 진행하도록 하겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 9: 마무리 (마무리/감사 1장 형식 개편)
    # 가이드: Q&A 안내, 연락처
    # -------------------------------------------------------------
    slide9 = prs.slides[8]
    clean_slide(slide9)
    # 헤더 없이 전면형으로 구성
    add_footer_and_pagenum(slide9, 9)
    
    center_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf9 = center_box.text_frame
    tf9.word_wrap = True
    
    p9_title = tf9.paragraphs[0]
    p9_title.text = "경청해 주셔서 감사합니다."
    p9_title.font.name = '맑은 고딕'
    p9_title.font.size = Pt(40)
    p9_title.font.bold = True
    p9_title.font.color.rgb = COLOR_PRIMARY
    p9_title.alignment = PP_ALIGN.CENTER
    p9_title.space_after = Pt(24)
    
    p9_sub = tf9.add_paragraph()
    p9_sub.text = "[ Q&A 및 자유 토론 ]"
    p9_sub.font.name = '맑은 고딕'
    p9_sub.font.size = Pt(20)
    p9_sub.font.bold = True
    p9_sub.font.color.rgb = COLOR_SECONDARY
    p9_sub.alignment = PP_ALIGN.CENTER
    p9_sub.space_after = Pt(24)
    
    p9_info = tf9.add_paragraph()
    p9_info.text = "문의처: 영업기획팀 김민서 대리 (minseo.kim@company.com / 내선번호 342)"
    p9_info.font.name = '맑은 고딕'
    p9_info.font.size = Pt(14)
    p9_info.font.color.rgb = COLOR_DARK
    p9_info.alignment = PP_ALIGN.CENTER
    
    add_notes(slide9, 
        "이상으로 2025년 상반기 매출 보고 발표를 모두 마치겠습니다. "
        "오늘 건의드린 하반기 구체적 대안들에 대해 질문이 있으시거나 의견이 있으시다면 질의하여 주시기 바랍니다. "
        "감사합니다.")
    
    prs.save(pptx_path)
    print(f"Restructured presentation successfully saved to: {pptx_path}")

if __name__ == "__main__":
    main()
