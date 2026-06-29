import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR

# UTF-8 출력 재설정 (윈도우 한글 인코딩 깨짐 방지)
sys.stdout.reconfigure(encoding='utf-8')

# 색상 상수 정의 (가이드 준수 및 고품질 디자인 적용)
COLOR_BG = RGBColor(255, 255, 255)         # 배경 흰색
COLOR_PRIMARY = RGBColor(43, 87, 151)       # 파랑 (#2B5797)
COLOR_SECONDARY = RGBColor(232, 119, 34)    # 포인트 주황 (#E87722)
COLOR_DARK = RGBColor(43, 87, 151)          # 메인 글씨 색상 (파랑 통일)
COLOR_MUTED = RGBColor(120, 130, 145)       # 보조 그레이 색상 (푸터/쪽번호용)
COLOR_CARD_BG = RGBColor(245, 247, 250)     # 카드 배경용 연한 회색 (Off-white)
COLOR_WHITE = RGBColor(255, 255, 255)

def clean_slide(slide):
    """슬라이드의 모든 기존 요소를 삭제하고 배경을 흰색으로 초기화"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG
    
    shapes = list(slide.shapes)
    for s in shapes:
        sp = s._element
        sp.getparent().remove(sp)

def add_header(slide, title_text, subtitle_text):
    """차별화된 헤더 추가 (데코 라인 하단 가로 주황선 배치 + 타이틀)"""
    # 타이틀 텍스트박스
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = '맑은 고딕'
    p.font.size = Pt(28)  # 28pt 이상, 굵게
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 서브 메시지 텍스트박스 (주황색 포인트 활용)
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.92), Inches(11.7), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_top = tf_sub.margin_bottom = tf_sub.margin_left = tf_sub.margin_right = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.name = '맑은 고딕'
    p_sub.font.size = Pt(15)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_SECONDARY
    
    # 하단 얇은 주황색 가로선 구분선 (이전과 디자인 다름)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_SECONDARY
    line.line.fill.background()

def add_footer_and_pagenum(slide, page_num):
    """슬라이드 하단에 부서명("전략기획팀") 및 쪽번호 추가 (2cm 여백 고려: 6.9 inches)"""
    # 부서명 (좌측)
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(4.0), Inches(0.3))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "전략기획팀"
    p_l.font.name = '맑은 고딕'
    p_l.font.size = Pt(11)
    p_l.font.color.rgb = COLOR_MUTED
    
    # 쪽번호 (우측)
    right_box = slide.shapes.add_textbox(Inches(11.0), Inches(6.9), Inches(1.5), Inches(0.3))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = str(page_num)
    p_r.font.name = '맑은 고딕'
    p_r.font.size = Pt(11)
    p_r.font.color.rgb = COLOR_MUTED
    p_r.alignment = PP_ALIGN.RIGHT

def add_accent_card(slide, left, top, width, height, accent_color=COLOR_PRIMARY, bg_color=COLOR_CARD_BG):
    """차별화된 디자인: 상단에 컬러 액센트 바가 있는 보더리스 플랫 카드 레이아웃"""
    # 1. 메인 카드 본체
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    base.fill.solid()
    base.fill.fore_color.rgb = bg_color
    base.line.fill.background()
    
    # 2. 상단 액센트 라인 (두께 4pt)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    return base

def format_text_box(slide, text_lines, left, top, width, height, font_size=18, bg_color=None, accent_color=None):
    """3줄 이하 텍스트 본문 생성 상자 (액센트 바 카드 옵션 지원)"""
    if bg_color and accent_color:
        add_accent_card(slide, left, top, width, height, accent_color, bg_color)
    elif bg_color:
        base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        base.fill.solid()
        base.fill.fore_color.rgb = bg_color
        base.line.fill.background()
        
    t_box = slide.shapes.add_textbox(left + Inches(0.2) if bg_color else left, 
                                     top + Inches(0.2) if bg_color else top, 
                                     width - Inches(0.4) if bg_color else width, 
                                     height - Inches(0.3) if bg_color else height)
    tf = t_box.text_frame
    tf.word_wrap = True
    
    for idx, line in enumerate(text_lines):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = line
        p.font.name = '맑은 고딕'
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_DARK
        p.font.bold = False
        p.space_after = Pt(12)

def add_notes(slide, notes_text):
    """슬라이드 노트 등록"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def main():
    prs = Presentation()
    # 16:9 와이드스크린으로 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide_layout = prs.slide_layouts[6] # 빈 슬라이드
    
    # -------------------------------------------------------------
    # SLIDE 1: 표지 (디자인 변경: 미니멀리즘 센터 정렬 레이아웃)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(slide_layout)
    clean_slide(slide1)
    
    # 상단 장식바
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.fill.background()
    
    # 중앙 타이틀
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(1.5))
    tf1 = title_box.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "2025년 하반기 사업 계획"
    p1.font.name = '맑은 고딕'
    p1.font.size = Pt(42)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    p1.alignment = PP_ALIGN.CENTER
    
    # 중앙 주황색 데코라인 (이전 표지와 다름)
    mid_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.166), Inches(3.6), Inches(3.0), Inches(0.04))
    mid_line.fill.solid()
    mid_line.fill.fore_color.rgb = COLOR_SECONDARY
    mid_line.line.fill.background()
    
    # 중앙 부제
    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(11.333), Inches(1.0))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "전략 기획 및 하반기 매출 32억 원 목표 달성안"
    p_sub.font.name = '맑은 고딕'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_SECONDARY
    p_sub.alignment = PP_ALIGN.CENTER
    
    # 하단 작성자 정보
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.0))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "보고대상 : 대표이사  |  작성부서 : 전략기획팀  |  작성일 : 2025년 7월 10일"
    p_info.font.name = '맑은 고딕'
    p_info.font.size = Pt(13)
    p_info.font.color.rgb = COLOR_MUTED
    p_info.alignment = PP_ALIGN.CENTER
    
    add_notes(slide1, 
        "안녕하십니까. 전략기획팀 보고를 맡은 김민서 대리입니다. 지금부터 2025년 하반기 사업 계획 발표를 시작하겠습니다. "
        "본 계획서는 상반기 실적을 토대로 하반기 매출 32억 원 목표를 성공적으로 달성하기 위한 구체적인 4대 전략과 세부 일정을 제시하고자 수립되었습니다. "
        "다음 슬라이드로 이동해 발표 목차를 안내해 드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 2: 목차 (디자인 변경: 4개 병합 슬림 항목 구성 및 격자형 카드 배치)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    clean_slide(slide2)
    add_header(slide2, "목차", "상반기 분석 결과를 활용한 하반기 핵심 4대 구성 순서입니다.")
    add_footer_and_pagenum(slide2, 2)
    
    toc_items = [
        ("01", "비전 및 하반기 목표", "경영 슬로건 및 하반기 주요 타겟 지표 요약"),
        ("02", "핵심 전략 4대 축", "프리미엄, B2B, 지역 개척, 고객 경험 전략"),
        ("03", "월별 로드맵 및 예상 손익", "하반기 일정 전개 계획 및 예상 재무 지표"),
        ("04", "리스크 및 대응 방안", "시장 리스크 예측 및 선제 조치 시나리오")
    ]
    
    col_w = Inches(5.6)
    row_h = Inches(1.8)
    start_x = Inches(0.8)
    start_y = Inches(1.8)
    
    for idx, (num, title, desc) in enumerate(toc_items):
        r = idx // 2
        c = idx % 2
        lx = start_x + c * (col_w + Inches(0.5))
        ty = start_y + r * (row_h + Inches(0.3))
        
        # 주황/파랑 번갈아가며 액센트 카드 생성
        accent_c = COLOR_PRIMARY if idx % 2 == 0 else COLOR_SECONDARY
        add_accent_card(slide2, lx, ty, col_w, row_h, accent_color=accent_c)
        
        # 번호 텍스트
        num_box = slide2.shapes.add_textbox(lx + Inches(0.2), ty + Inches(0.15), Inches(1.0), Inches(0.4))
        p_n = num_box.text_frame.paragraphs[0]
        p_n.text = num
        p_n.font.name = '맑은 고딕'
        p_n.font.size = Pt(20)
        p_n.font.bold = True
        p_n.font.color.rgb = accent_c
        
        # 타이틀 및 설명
        t_box = slide2.shapes.add_textbox(lx + Inches(1.0), ty + Inches(0.15), col_w - Inches(1.2), Inches(1.4))
        tf = t_box.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = '맑은 고딕'
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_DARK
        p_t.space_after = Pt(6)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = '맑은 고딕'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_MUTED
        
    add_notes(slide2, 
        "이 슬라이드에서 말씀드릴 것은 오늘 보고드릴 4가지 핵심 목차입니다. "
        "구체적으로 말씀드리면, 먼저 당사의 하반기 도약 비전과 지표 목표를 공유하고, "
        "이를 수행할 핵심 4대 사업 전략을 예산과 연계하여 보고하겠습니다. "
        "이어서 월별 상세 실행안 및 예상 이익, 그리고 발생 가능한 리스크 대응책 순으로 보고드리겠습니다. "
        "첫 번째 순서인 비전 및 목표부터 보고해 드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 3: 비전 및 하반기 목표 (3줄 이내 요약)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    clean_slide(slide3)
    add_header(slide3, "1. 비전 및 하반기 목표", "하반기 매출 32억 원 달성 및 고객 추천 만족도 NPS 78점 이상 확보")
    add_footer_and_pagenum(slide3, 3)
    
    targets = [
        "· 비전 슬로건: \"디지털 전환의 동반자, 고객 성공의 파트너\" 대내외 전파",
        "· 영업 핵심 지표: 하반기 매출 32억 원(+12.6% YoY) 및 신규 고객 350건 수주 달성",
        "· 품질 관리 지표: 고객 충성 만족도 NPS 78점 돌파 및 서비스 갱신율 극대화"
    ]
    format_text_box(slide3, targets, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), font_size=18, bg_color=COLOR_CARD_BG, accent_color=COLOR_PRIMARY)
    
    add_notes(slide3, 
        "이 슬라이드에서 말씀드릴 것은 당사의 하반기 지향 비전과 정량적 목표 요약입니다. "
        "구체적으로 말씀드리면, 당사는 '디지털 전환의 동반자, 고객 성공의 파트너'라는 새로운 비전 아래, "
        "하반기 매출 목표 32억 원을 수립했습니다. 이는 상반기 성과 대비 12.6% 증가한 수치입니다. "
        "또한 신규 고객 350건 유치와 품질 강화를 통한 NPS 78점 이상 확보를 핵심 지표로 설정했습니다. "
        "그럼 다음 슬라이드에서는 이러한 성과를 만들기 위한 핵심 4대 전략과 예산 비중을 살펴보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 4: 핵심 전략 4대 축 및 예산 비중 (원그래프 삽입)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    clean_slide(slide4)
    add_header(slide4, "2. 핵심 전략 4대 축", "전략적 자원 배분: 총 예산 6.5억 원의 63%를 프리미엄 및 B2B 확장에 집중")
    add_footer_and_pagenum(slide4, 4)
    
    # 원그래프 데이터 준비 (전략별 예산 억 원 단위)
    chart_data_pie = CategoryChartData()
    chart_data_pie.categories = [
        '프리미엄 라인 (2.3억)', 
        '엔터프라이즈 B2B (1.8억)', 
        '지역 시장 개척 (1.5억)', 
        '고객 경험 혁신 (0.9억)'
    ]
    chart_data_pie.add_series('예산 배분 (억 원)', (2.3, 1.8, 1.5, 0.9))
    
    # 원그래프 (Pie Chart) 삽입
    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.5)
    chart_pie = slide4.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data_pie
    ).chart
    chart_pie.has_legend = True
    chart_pie.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # 우측 텍스트 설명 카드 (액센트 적용)
    format_text_box(slide4, [
        "■ 4대 핵심 전략 요약",
        "· 프리미엄 가속: AI 보고 기능 탑재(V2.0) 및 업셀링",
        "· B2B 영업 활성화: 금융·공공팀 신설 및 세미나 정례화",
        "· 비수도권 개척: 부산·대구 파트너 센터 개설 및 지역 할인",
        "· 고객 혁신: 가입 3일 내 온보딩 개편 및 챗봇 가동"
    ], Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.5), font_size=13, bg_color=COLOR_CARD_BG, accent_color=COLOR_SECONDARY)
    
    add_notes(slide4, 
        "이 슬라이드에서 말씀드릴 것은 하반기 4대 핵심 전략과 사업 예산 비중입니다. "
        "구체적으로 말씀드리면, 전체 예산 6억 5,000만 원 중 프리미엄 라인 고도화에 2억 3천만 원, 엔터프라이즈 B2B 영업망 확장에 1억 8천만 원을 책정하여 상반기 핵심 수익 모델에 리소스를 집중 투자했습니다. "
        "아울러 지방 시장 개척에 1억 5천만 원, 고객 경험 혁신에 9천만 원을 안배하여 균형 있는 성장을 도모하고자 합니다. "
        "그럼 다음 슬라이드로 이동하여 월별 추진 일정을 로드맵으로 확인해 보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 5: 월별 로드맵 (디자인 변경: 2x3 레이아웃의 타임라인 카드 배치)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    clean_slide(slide5)
    add_header(slide5, "3. 월별 로드맵", "하반기 핵심 서비스 출시 및 B2B/지역 개척 일정 타임라인 전개")
    add_footer_and_pagenum(slide5, 5)
    
    roadmap_items = [
        ("7월", "프리미엄 V2.0 출시", "AI 자동 보고서 기능 런칭 및 킥오프"),
        ("8월", "B2B 전담 영업팀 신설", "금융·공공 타겟 채용 및 캠페인 가동"),
        ("9월", "부산·대구 지역 센터", "지방 거점 파트너 계약 및 개소"),
        ("10월", "교육기관 패키지 런칭", "24시간 챗봇 상담 개시 및 고객 혁신"),
        ("11월", "지자체 디지털 교육", "공공 부문 제안서 접수 및 마케팅"),
        ("12월", "성과 분석 및 정산", "차년도 영업 전략 및 예산안 확정")
    ]
    
    col_w = Inches(3.7)
    row_h = Inches(1.8)
    start_x = Inches(0.8)
    start_y = Inches(1.8)
    
    for idx, (month, title, desc) in enumerate(roadmap_items):
        r = idx // 3
        c = idx % 3
        lx = start_x + c * (col_w + Inches(0.3))
        ty = start_y + r * (row_h + Inches(0.3))
        
        # 짝수달은 파란색 액센트, 홀수달은 주황색 액센트
        col_acc = COLOR_PRIMARY if idx % 2 == 0 else COLOR_SECONDARY
        add_accent_card(slide5, lx, ty, col_w, row_h, accent_color=col_acc)
        
        # 텍스트 정보
        t_box = slide5.shapes.add_textbox(lx + Inches(0.15), ty + Inches(0.15), col_w - Inches(0.3), row_h - Inches(0.3))
        tf = t_box.text_frame
        tf.word_wrap = True
        
        p_m = tf.paragraphs[0]
        p_m.text = month
        p_m.font.name = '맑은 고딕'
        p_m.font.size = Pt(18)
        p_m.font.bold = True
        p_m.font.color.rgb = col_acc
        p_m.space_after = Pt(4)
        
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.name = '맑은 고딕'
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_DARK
        p_t.space_after = Pt(4)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = '맑은 고딕'
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_MUTED
        
    add_notes(slide5, 
        "이 슬라이드에서 말씀드릴 것은 하반기 추진 월별 로드맵입니다. "
        "구체적으로 말씀드리면, 7월에 프리미엄 V2.0 버전 출시를 마친 뒤, 8월 B2B 금융/공공팀 조직 신설을 단행하겠습니다. "
        "9월에는 부산과 대구에 지역 파트너 센터를 개소하여 전국 커버리지를 늘릴 것입니다. "
        "이후 10월 교육기관용 특화 패키지 런칭과 24시간 CS 인프라를 가동하고, 11월 지자체 제안을 추진하여 12월 최종 실적을 결산하겠습니다. "
        "이러한 로드맵 실행을 통한 예상 재무 손익을 이어서 보고드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 6: 예상 손익 요약 (세로 막대 차트 삽입)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    clean_slide(slide6)
    add_header(slide6, "4. 예상 손익", "하반기 매출 32.0억 원 및 비용 22.4억 원 집행, 영업이익률 30.0% 마크 목표")
    add_footer_and_pagenum(slide6, 6)
    
    # 세로 막대 차트 데이터 준비 (비교 데이터)
    chart_data_col = CategoryChartData()
    chart_data_col.categories = ['예상 매출', '예상 비용', '예상 영업이익']
    chart_data_col.add_series('금액 (억 원)', (32.0, 22.4, 9.6))
    
    # Clustered Column 차트 삽입
    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.5)
    chart_col = slide6.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_col
    ).chart
    chart_col.has_legend = False
    
    # 우측 세부 비용 테이블 성격 텍스트 박스
    format_text_box(slide6, [
        "■ 예상 비용 상세 집행액",
        "· 인건비: 14억 2,000만 원 (핵심 인력 3명 포함)",
        "· 마케팅비: 4억 6,000만 원 (대규모 캠페인 전개)",
        "· 운영비: 3억 6,000만 원 (지사 센터 및 인프라)",
        "※ 영업이익률: 30.0% 달성 전망"
    ], Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.5), font_size=13, bg_color=COLOR_CARD_BG, accent_color=COLOR_PRIMARY)
    
    add_notes(slide6, 
        "이 슬라이드에서 말씀드릴 것은 하반기 예상 손익 전망입니다. "
        "구체적으로 말씀드리면, 매출 32억 원 달성 대비 예상 비용은 총 22억 4천만 원이 예상됩니다. "
        "비용의 상세 항목은 R&D 인건비 14억 2천만 원, 프리미엄 캠페인용 마케팅비 4억 6천만 원, 그리고 운영비 3억 6천만 원으로 계획했습니다. "
        "이를 통해 영업이익 9억 6천만 원을 수확하고 30.0%의 우수한 영업이익률을 달성하겠습니다. "
        "그럼 다음 슬라이드에서 극복할 잠재 리스크와 그 대응책을 말씀드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 7: 리스크 및 대응 방안 (3줄 이내 요약)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    clean_slide(slide7)
    add_header(slide7, "5. 리스크 및 대응 방안", "경기 둔화, 경쟁사 공격 및 인력 확보 등 핵심 3대 위협의 완벽한 리스크 헤징")
    add_footer_and_pagenum(slide7, 7)
    
    risks = [
        "· 1) 경기 둔화 위험: 솔루션 분할 결제 등 유연 요금제 선제적 도입 (진입 장벽 완화)",
        "· 2) 경쟁사 신제품 출현: AI 핵심 기능을 탑재한 프리미엄 V2.0 출시 앞당김 (9월 -> 8월)",
        "· 3) 인력 이탈 및 공백: 7~8월 내 마케팅 및 개발 핵심 부문 3명 특별 특별 전담 채용 완료"
    ]
    format_text_box(slide7, risks, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), font_size=18, bg_color=COLOR_CARD_BG, accent_color=COLOR_SECONDARY)
    
    add_notes(slide7, 
        "이 슬라이드에서 말씀드릴 것은 예상되는 위협 리스크 및 이에 대한 예방 조치입니다. "
        "구체적으로 말씀드리면, 첫째, 거시 경기 둔화 우려에 대처하기 위해 장기 분할 납부를 도입하겠습니다. "
        "둘째, 타사의 유사 솔루션 견제를 위해 핵심 AI 로드맵 일정을 9월에서 8월로 한 달 당겨 조기 런칭하겠습니다. "
        "셋째, 일시적 업무 병목을 막고자 8월까지 마케팅과 개발 인원 3명을 충원하겠습니다. "
        "이어서 하반기 사업 계획의 종합 요약과 제안 단계를 최종 보고해 드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 8: 요약 및 제안 (요약/제안 1장 형식 준수)
    # 가이드: 핵심 발견 3줄 / 제안, 다음 단계 2줄
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(slide_layout)
    clean_slide(slide8)
    add_header(slide8, "요약 및 제안", "하반기 매출 32억 원과 영업이익 9.6억 원 달성을 위한 제언 요약")
    add_footer_and_pagenum(slide8, 8)
    
    # 핵심 발견 3줄 박스 (주황 액센트 적용)
    format_text_box(slide8, [
        "[핵심 발견 (Key Findings)]",
        "1. 매출 32억 원 목표를 위해 상반기 성장 모멘텀이 증명된 프리미엄 및 B2B 영업 가동 필요",
        "2. 전국 영업 기반 구축 및 지방 잠재 수요 흡수를 위해 비수도권 매출 비중 45% 유치 중요",
        "3. 프리미엄 신기능 V2.0 조기 출시(8월)를 통한 시장 경쟁력 선점 및 리스크 관리 요구됨"
    ], Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2), font_size=13, bg_color=COLOR_CARD_BG, accent_color=COLOR_SECONDARY)
    
    # 제안/다음 단계 2줄 박스 (파랑 액센트 적용)
    format_text_box(slide8, [
        "[제안 및 다음 단계 (Proposals & Next Steps)]",
        "1. 대표이사 승인 직후 전 부서 대상 하반기 킥오프 세션 진행 및 마일스톤 실행 지시",
        "2. 7월 중 신규 채용 승인 가결 및 프리미엄 V2.0 AI 보고 모듈 개발 진척도 우선 감시"
    ], Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.0), font_size=13, bg_color=COLOR_CARD_BG, accent_color=COLOR_PRIMARY)
    
    add_notes(slide8, 
        "이 슬라이드에서 말씀드릴 것은 하반기 비즈니스의 종합 제언입니다. "
        "구체적으로 말씀드리면, 하반기 매출 32억 원 고지를 탈환하기 위해선 프리미엄 요금 업셀링과 신설 금융/공공팀의 정밀 영업이 결합되어야 합니다. "
        "아울러 지방 시장 활성화와 AI 출시 마일스톤 준수가 선행 과제입니다. "
        "이에 대한 건의 사항으로 승인 통과 즉시 킥오프 워크숍을 가동하고 7월 R&D 일정 정밀 진단을 요구드립니다. "
        "그럼 마지막 Q&A 시간을 갖겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 9: 마무리 및 Q&A (마무리/감사 1장 형식 준수)
    # 가이드: Q&A 안내, 연락처
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(slide_layout)
    clean_slide(slide9)
    add_footer_and_pagenum(slide9, 9)
    
    center_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf9 = center_box.text_frame
    tf9.word_wrap = True
    
    p9_title = tf9.paragraphs[0]
    p9_title.text = "경청해 주셔서 감사합니다."
    p9_title.font.name = '맑은 고딕'
    p9_title.font.size = Pt(40)
    p9_title.font.bold = True
    p9_title.font.color.rgb = COLOR_PRIMARY
    p9_title.alignment = PP_ALIGN.CENTER
    p9_title.space_after = Pt(24)
    
    p9_sub = tf9.add_paragraph()
    p9_sub.text = "[ Q&A 및 토론 ]"
    p9_sub.font.name = '맑은 고딕'
    p9_sub.font.size = Pt(20)
    p9_sub.font.bold = True
    p9_sub.font.color.rgb = COLOR_SECONDARY
    p9_sub.alignment = PP_ALIGN.CENTER
    p9_sub.space_after = Pt(24)
    
    p9_info = tf9.add_paragraph()
    p9_info.text = "문의처: 전략기획팀 (strategy@company.com / 사내전화 101)"
    p9_info.font.name = '맑은 고딕'
    p9_info.font.size = Pt(14)
    p9_info.font.color.rgb = COLOR_DARK
    p9_info.alignment = PP_ALIGN.CENTER
    
    add_notes(slide9, 
        "이상으로 2025년 하반기 사업 계획 발표를 모두 마치고 Q&A로 전환하겠습니다. "
        "오늘 공유드린 전략 과제와 손익 전망에 대해 건의사항이나 질의가 있으시다면 편안하게 개진해 주시기 바랍니다. "
        "감사합니다.")
    
    # 저장 경로
    output_path = "c:/Users/user/Desktop/ai사무자동화/예제파일/12_발표_프롬프트_실습/02_발표_재료/하반기_사업계획_발표.pptx"
    prs.save(output_path)
    print(f"Business plan presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
