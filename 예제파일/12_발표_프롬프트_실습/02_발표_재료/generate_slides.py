import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR

# 색상 상수 정의 (가이드 준수 및 고품질 디자인 적용)
COLOR_BG = RGBColor(248, 249, 250)         # 연한 그레이 (배경)
COLOR_PRIMARY = RGBColor(43, 87, 151)       # 신뢰감을 주는 파랑 (#2B5797)
COLOR_SECONDARY = RGBColor(232, 119, 34)    # 포인트 주황 (#E87722)
COLOR_DARK = RGBColor(33, 37, 41)           # 메인 텍스트 (다크 그레이)
COLOR_MUTED = RGBColor(108, 117, 125)       # 보조 텍스트 (그레이)
COLOR_WHITE = RGBColor(255, 255, 255)       # 카드 배경 (화이트)
COLOR_LIGHT_BLUE = RGBColor(230, 240, 255)  # 서브 하이라이트 배경
COLOR_LIGHT_ORANGE = RGBColor(254, 243, 235)# 서브 하이라이트 주황 배경

def set_slide_background(slide):
    """슬라이드 배경색을 연한 그레이(#F8F9FA)로 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, subtitle_text):
    """표준화된 세련된 헤더 영역 추가 (좌측 데코 바 + 타이틀 + 키 메시지 + 구분선)"""
    # 1. 좌측 파란색 포인트 바
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(0.12), Inches(0.8)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_PRIMARY
    left_bar.line.fill.background() # 선 없음
    
    # 2. 타이틀 텍스트박스
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.42), Inches(11.5), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = '맑은 고딕'
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 3. 키 메시지 텍스트박스 (주황 포인트 컬러)
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.92), Inches(11.5), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.name = '맑은 고딕'
    p_sub.font.size = Pt(14)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_SECONDARY
    
    # 4. 하단 얇은 구분선
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 224, 230)
    line.line.fill.background()

def create_card(slide, left, top, width, height, bg_color=COLOR_WHITE):
    """깔끔한 화이트/컬러 카드 레이아웃 컨테이너 생성"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    # 테두리를 매우 연하게 처리하여 모던한 플랫 디자인 연출
    card.line.color.rgb = RGBColor(230, 235, 240)
    card.line.width = Pt(1.5)
    return card

def add_notes(slide, notes_text):
    """슬라이드 노트에 발표 대본 기입"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def format_para(paragraph, font_name='맑은 고딕', size_pt=14, bold=False, color_rgb=COLOR_DARK, align=PP_ALIGN.LEFT):
    """문단 서식 설정 헬퍼"""
    paragraph.font.name = font_name
    paragraph.font.size = Pt(size_pt)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color_rgb
    paragraph.alignment = align

def main():
    csv_path = "c:/Users/user/Desktop/ai사무자동화/예제파일/12_발표_프롬프트_실습/02_발표_재료/월별매출.csv"
    
    # 데이터 분석
    df = pd.read_csv(csv_path)
    # 매출액 단위 변환 (CSV 값은 100원 단위이므로 10만으로 나누면 '억 원' 단위, 10으로 나누면 '만 원' 단위)
    # 예: 3,168,000 / 1,000,000 = 3.168억 원
    df['매출액_억원'] = df['매출액'] / 1000000.0
    df['매출액_만원'] = df['매출액'] / 10.0
    
    # 월별/제품별 피벗
    pivot_sales = df.pivot(index='월', columns='제품', values='매출액_억원')
    # 월 정렬을 위해 순서 정의
    months_order = ['1월', '2월', '3월', '4월', '5월', '6월']
    pivot_sales = pivot_sales.reindex(months_order)
    
    prs = Presentation()
    # 16:9 와이드스크린 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # SLIDE 1: 표지 (Cover) - 다크 네이비 테마로 신뢰도 높고 고급스러운 느낌 연출
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # 빈 슬라이드
    slide1 = prs.slides.add_slide(slide_layout)
    
    # 다크 배경
    bg_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLOR_PRIMARY
    bg_shape.line.fill.background()
    
    # 데코 포인트 (오렌지 띠)
    deco_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))
    deco_shape.fill.solid()
    deco_shape.fill.fore_color.rgb = COLOR_SECONDARY
    deco_shape.line.fill.background()
    
    # 메인 타이틀 박스
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "2025년 상반기 매출 보고 요약"
    format_para(p1, size_pt=40, bold=True, color_rgb=COLOR_WHITE)
    
    # 부제 박스
    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.3), Inches(1.0))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "보고 기간: 2025년 1월 ~ 6월  |  보고 대상: 경영지원본부장"
    format_para(p_sub, size_pt=18, bold=False, color_rgb=RGBColor(210, 225, 250))
    
    # 작성자 정보 박스
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.0))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "작성부서 : 영업기획팀\n작 성 자 : 김민서 대리\n작 성 일 : 2025년 7월 5일"
    format_para(p_info, size_pt=14, bold=False, color_rgb=RGBColor(230, 240, 255))
    p_info.line_spacing = 1.3
    
    add_notes(slide1, 
        "안녕하십니까. 영업기획팀 김민서 대리입니다. 지금부터 2025년 상반기 매출 보고 요약 발표를 시작하도록 하겠습니다. "
        "이번 발표는 지난 1월부터 6월까지의 상반기 영업 실적을 체계적으로 분석하여 핵심 성과와 이슈를 공유하고, "
        "이를 극복하기 위한 하반기 전략 제안을 드리고자 준비했습니다. 그럼 다음 슬라이드로 이동해 목차부터 살펴보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 2: 목차 (Table of Contents)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "목차", "상반기 실적 분석부터 하반기 핵심 전략 제안까지의 진행 흐름입니다.")
    
    # 목차 아이템 배치 (2열로 배치하여 여백 확보 및 모던한 느낌 제공)
    toc_items = [
        ("01", "전체 실적 요약", "상반기 매출/판매건수 등 총괄 실적 종합"),
        ("02", "제품별 분석", "베이직, 프리미엄, 엔터프라이즈별 실적 분석"),
        ("03", "분기별 매출 추이", "1분기 대비 2분기 주요 성장 요인 진단"),
        ("04", "지역별 매출 분포", "수도권 편중 현상 및 지역별 실적 요약"),
        ("05", "주요 성과 및 이슈", "신규 유치 목표 달성과 비수기 공백 등 진단"),
        ("06", "하반기 전략 제안", "성장세를 이어가기 위한 4대 액션 플랜")
    ]
    
    col_w = Inches(5.6)
    row_h = Inches(1.4)
    start_x = Inches(0.8)
    start_y = Inches(1.8)
    
    for i, (num, title, desc) in enumerate(toc_items):
        r = i // 2
        c = i % 2
        lx = start_x + c * (col_w + Inches(0.5))
        ty = start_y + r * (row_h + Inches(0.3))
        
        # 각 목차 카드 생성
        create_card(slide2, lx, ty, col_w, row_h)
        
        # 숫자 아이콘 생성
        num_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx + Inches(0.2), ty + Inches(0.2), Inches(0.8), Inches(0.8))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_LIGHT_BLUE
        num_box.line.fill.background()
        tf_num = num_box.text_frame
        tf_num.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_num = tf_num.paragraphs[0]
        p_num.text = num
        format_para(p_num, size_pt=20, bold=True, color_rgb=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
        
        # 텍스트 정보
        text_box = slide2.shapes.add_textbox(lx + Inches(1.2), ty + Inches(0.15), col_w - Inches(1.4), Inches(1.0))
        tf_text = text_box.text_frame
        tf_text.word_wrap = True
        p_title = tf_text.paragraphs[0]
        p_title.text = title
        format_para(p_title, size_pt=16, bold=True, color_rgb=COLOR_DARK)
        
        p_desc = tf_text.add_paragraph()
        p_desc.text = desc
        format_para(p_desc, size_pt=12, bold=False, color_rgb=COLOR_MUTED)
        
    add_notes(slide2, 
        "이 슬라이드에서 말씀드릴 것은 오늘 발표할 전체 흐름입니다. "
        "발표는 전체 실적 요약을 시작으로, 제품군별 세부 매출 비중, 분기별 매출 추이, 그리고 지역별 매출 분포를 다루겠습니다. "
        "이어서 상반기 활동을 총평하며 주요 성과와 극복 과제를 분석하고, 마지막으로 하반기 구체적인 전략 제안 순으로 진행하겠습니다. "
        "그럼 첫 번째 순서인 전체 실적 요약부터 보고해 드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 3: 전체 실적 요약 (Overview)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "1. 전체 실적 요약", "상반기 총매출 28억 4,200만 원 달성, 전년 동기 대비 +12.3% 성장 기록")
    
    # 4개 핵심 지표 카드 배치
    kpis = [
        ("상반기 총매출", "28억 4,200만 원", "전년 동기 대비 +12.3%", COLOR_PRIMARY),
        ("총 판매 건수", "1,847 건", "전년 동기 대비 +8.7%", COLOR_PRIMARY),
        ("월평균 매출", "4억 7,367만 원", "상반기 평균 실적", COLOR_DARK),
        ("매출 극단값", "최고 4월 / 최저 3월", "4월: 6.9억 / 3월: 2.08억", COLOR_SECONDARY)
    ]
    
    card_w = Inches(2.7)
    card_h = Inches(4.5)
    start_x = Inches(0.8)
    spacing = Inches(0.3)
    
    for i, (title, val, desc, color) in enumerate(kpis):
        lx = start_x + i * (card_w + spacing)
        # 카드 배경
        create_card(slide3, lx, Inches(1.8), card_w, card_h)
        
        # 타이틀 영역 (배경바 포함)
        header_bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, Inches(1.8), card_w, Inches(0.8))
        header_bar.fill.solid()
        # 마지막 극단값 카드는 오렌지 테마로 하이라이트
        if color == COLOR_SECONDARY:
            header_bar.fill.fore_color.rgb = COLOR_LIGHT_ORANGE
            title_color = COLOR_SECONDARY
        else:
            header_bar.fill.fore_color.rgb = COLOR_LIGHT_BLUE
            title_color = COLOR_PRIMARY
            
        header_bar.line.fill.background()
        
        tf_h = header_bar.text_frame
        tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        format_para(p_h, size_pt=15, bold=True, color_rgb=title_color, align=PP_ALIGN.CENTER)
        
        # 값 & 설명 텍스트
        val_box = slide3.shapes.add_textbox(lx + Inches(0.1), Inches(2.9), card_w - Inches(0.2), Inches(3.0))
        tf_val = val_box.text_frame
        tf_val.word_wrap = True
        
        p_v = tf_val.paragraphs[0]
        p_v.text = val
        # 글자 크기가 작게 들어가지 않도록 중요 값은 24pt 이상 설정
        format_para(p_v, size_pt=24, bold=True, color_rgb=color, align=PP_ALIGN.CENTER)
        p_v.space_after = Pt(20)
        
        p_d = tf_val.add_paragraph()
        p_d.text = desc
        format_para(p_d, size_pt=13, bold=False, color_rgb=COLOR_MUTED, align=PP_ALIGN.CENTER)
        
    add_notes(slide3, 
        "이 슬라이드에서 말씀드릴 것은 2025년 상반기 전체 실적의 핵심 요약입니다. "
        "구체적으로 말씀드리면, 상반기 총매출은 28억 4,200만 원으로 전년 동기 대비 12.3% 큰 폭의 성장을 기록했습니다. "
        "총 판매 건수 또한 1,847건으로 전년 동기 대비 8.7% 증가하며 성장세를 보였습니다. "
        "월평균 매출은 4억 7,367만 원 수준입니다. "
        "특히 주목할 만한 부분은 매출의 변동폭인데, 봄맞이 프로모션 효과를 톡톡히 본 4월이 6억 9,003만 원으로 최고 매출을 달성한 반면, "
        "분기 전환기 비수기였던 3월은 2억 790만 원으로 최저 매출을 보였습니다. "
        "이러한 극단적인 실적 차이는 비수기 매출 보완 전략의 필요성을 시사합니다. "
        "그럼 다음 슬라이드에서는 이러한 실적을 견인한 제품군별 상세 데이터를 분석해 보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 4: 제품별 분석 (Product Analysis)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "2. 제품별 분석", "프리미엄(44.2%)과 엔터프라이즈(39.1%)가 상반기 매출 성장의 주도적 역할 수행")
    
    # 3개 제품군 카드 배치
    products = [
        ("베이직 (Basic)", "4억 7,322만 원", "16.7%", 
         ["· 안정적인 기반 매출 유지", "· 소규모 사업자 신규 유입 증가", "· 전년 대비: +5.2% 성장", "· 이탈률: 4.2% (전년 5.8% 대비 개선)"],
         COLOR_MUTED, COLOR_WHITE),
        ("프리미엄 (Premium)", "12억 5,580만 원", "44.2%", 
         ["· 2분기 뚜렷한 성장세 (4~6월 평균 19.8%↑)", "· 기능 업그레이드 패키지 인기", "· 전년 대비: +18.6% (최고 성장률)", "· 업셀링 성공률: 23.4% 기록"],
         COLOR_SECONDARY, COLOR_LIGHT_ORANGE), # 주황 하이라이트 배경
        ("엔터프라이즈 (Enterprise)", "11억 1,298만 원", "39.1%", 
         ["· 대형 고객사 3개사 신규 계약", "· B2B 영업팀 확충 효과 본격화", "· 전년 대비: +14.1% 성장", "· 계약 갱신율: 91.2% (업계평균 82% 상회)"],
         COLOR_PRIMARY, COLOR_LIGHT_BLUE) # 파랑 하이라이트 배경
    ]
    
    card_w = Inches(3.7)
    card_h = Inches(4.5)
    start_x = Inches(0.8)
    spacing = Inches(0.3)
    
    for i, (name, revenue, ratio, details, title_color, bg_col) in enumerate(products):
        lx = start_x + i * (card_w + spacing)
        # 카드 컨테이너 생성
        create_card(slide4, lx, Inches(1.8), card_w, card_h, bg_color=bg_col)
        
        # 텍스트 박스
        t_box = slide4.shapes.add_textbox(lx + Inches(0.2), Inches(2.0), card_w - Inches(0.4), card_h - Inches(0.4))
        tf = t_box.text_frame
        tf.word_wrap = True
        
        # 제품명
        p_name = tf.paragraphs[0]
        p_name.text = name
        format_para(p_name, size_pt=20, bold=True, color_rgb=title_color)
        p_name.space_after = Pt(10)
        
        # 매출액 & 비중
        p_rev = tf.add_paragraph()
        p_rev.text = f"{revenue} ({ratio})"
        format_para(p_rev, size_pt=22, bold=True, color_rgb=COLOR_DARK)
        p_rev.space_after = Pt(20)
        
        # 구분선 역할의 빈 줄
        p_line = tf.add_paragraph()
        p_line.text = "──────────────────"
        format_para(p_line, size_pt=10, bold=False, color_rgb=RGBColor(200, 205, 210))
        p_line.space_after = Pt(15)
        
        # 세부 항목
        for item in details:
            p_det = tf.add_paragraph()
            p_det.text = item
            format_para(p_det, size_pt=13, bold=False, color_rgb=COLOR_DARK)
            p_det.space_after = Pt(8)
            
    add_notes(slide4, 
        "이 슬라이드에서 말씀드릴 것은 제품군별 매출 분포 및 성과 분석입니다. "
        "구체적으로 말씀드리면, 매출 기여도가 가장 높은 제품군은 '프리미엄'으로, 12억 5,580만 원의 실적을 기록하며 전체 매출의 44.2%를 견인했습니다. "
        "프리미엄은 기능 업그레이드 패키지의 인기에 힘입어 전년 동기 대비 18.6% 성장하였으며, 베이직에서 프리미엄으로의 업셀링 성공률도 23.4%로 높게 나타났습니다. "
        "두 번째로 비중이 높은 '엔터프라이즈'는 11억 1,298만 원으로 전체의 39.1%를 차지했습니다. B2B 영업팀 확충 효과로 대형 신규 계약 3건을 수주했으며, 계약 갱신율도 91.2%로 업계 평균인 82%를 크게 상회했습니다. "
        "마지막으로 '베이직'은 4억 7,322만 원으로 전체의 16.7%를 차지하였으며, 소규모 사업자 신규 유입 덕에 전년비 5.2% 성장하며 안정적인 기반 매출 역할을 수행하고 있습니다. "
        "종합해 보면, 프리미엄과 엔터프라이즈가 전체 매출의 83% 이상을 책임지며 주 성장 동력 역할을 수행하고 있음을 알 수 있습니다. "
        "그럼 이어서 상반기 월별 및 분기별 매출 흐름을 꺾은선 그래프로 확인해 보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 5: 월별/분기별 추이 (Monthly & Quarterly Trends)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "3. 분기별 매출 추이", "1분기 대비 2분기 매출 +22.3% 급성장 (4월 프로모션 및 5월 대형 사업 수주)")
    
    # 좌측 영역: 차트 데이터 생성 및 꺾은선 차트 삽입
    chart_data = CategoryChartData()
    chart_data.categories = months_order
    # 제품별 데이터를 억 원 단위로 추가
    basic_data = [df[(df['월'] == m) & (df['제품'] == '베이직')]['매출액_억원'].values[0] for m in months_order]
    premium_data = [df[(df['월'] == m) & (df['제품'] == '프리미엄')]['매출액_억원'].values[0] for m in months_order]
    enterprise_data = [df[(df['월'] == m) & (df['제품'] == '엔터프라이즈')]['매출액_억원'].values[0] for m in months_order]
    
    chart_data.add_series('베이직', tuple(basic_data))
    chart_data.add_series('프리미엄', tuple(premium_data))
    chart_data.add_series('엔터프라이즈', tuple(enterprise_data))
    
    # 전체 월별 합산 라인을 그리기 위해 월별 총액 계산
    monthly_totals = []
    for m in months_order:
        tot = df[df['월']==m]['매출액_억원'].sum()
        monthly_totals.append(tot)
    chart_data.add_series('월 합계 매출', tuple(monthly_totals))
    
    # 차트 위치 및 크기
    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.8)
    chart = slide5.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    # 우측 영역: 분석 텍스트 카드
    create_card(slide5, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.8))
    
    t_box = slide5.shapes.add_textbox(Inches(8.2), Inches(2.0), Inches(4.1), Inches(4.4))
    tf = t_box.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = "분기별 주요 실적 요약"
    format_para(p_title, size_pt=18, bold=True, color_rgb=COLOR_PRIMARY)
    p_title.space_after = Pt(15)
    
    p_q1 = tf.add_paragraph()
    p_q1.text = "■ 1분기 (1~3월): 12억 7,830만 원"
    format_para(p_q1, size_pt=14, bold=True, color_rgb=COLOR_DARK)
    
    p_q1_desc = tf.add_paragraph()
    p_q1_desc.text = "  · 전년 동기 대비 +9.1% 성장\n  · 1월 명절 효과로 하락했으나 2월 빠르게 반등\n  · 3월은 계절성 비수기로 최저점(2억 790만 원) 기록"
    format_para(p_q1_desc, size_pt=12, bold=False, color_rgb=COLOR_MUTED)
    p_q1_desc.space_after = Pt(15)
    
    p_q2 = tf.add_paragraph()
    p_q2.text = "■ 2분기 (4~6월): 15억 6,370만 원"
    format_para(p_q2, size_pt=14, bold=True, color_rgb=COLOR_SECONDARY)
    
    p_q2_desc = tf.add_paragraph()
    p_q2_desc.text = "  · 전년 동기 대비 +15.8% 대폭 성장\n  · 4월 봄맞이 프로모션 성공으로 매출 최고점 돌파\n  · 5월 정부 디지털 전환 사업 수주 효과 반영"
    format_para(p_q2_desc, size_pt=12, bold=False, color_rgb=COLOR_MUTED)
    
    add_notes(slide5, 
        "이 슬라이드에서 말씀드릴 것은 월별 상세 매출 추이와 분기별 주요 변동 사유입니다. "
        "구체적으로 말씀드리면, 1분기 총 매출은 12억 7,830만 원을 기록하여 전년 대비 9.1% 성장했습니다. "
        "1월 명절 연휴 기간 중 일시적인 판매 하락이 있었으나 2월에 바로 회복되었습니다. 단, 3월은 계절적인 비수기가 맞물리면서 상반기 최저점인 2억 790만 원에 그쳤습니다. "
        "반면 2분기는 15억 6,370만 원을 달성하며 전년 동기 대비 15.8%라는 높은 성장률을 나타냈습니다. "
        "특히 4월에는 프리미엄/엔터프라이즈의 봄맞이 프로모션 성공으로 매출이 6억 9,003만 원까지 치솟았으며, "
        "5월에는 대형 공공기관이 포함된 정부 디지털 전환 사업 수주가 발생하면서 실적을 굳건히 지탱했습니다. "
        "이어서 다음 슬라이드에서는 지역별 매출의 상세 분포를 통해 영업 전략의 침투 수준을 점검하겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 6: 지역별 매출 분포 (Regional Distribution)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "4. 지역별 매출 분포", "수도권 매출 비중 62.3% 편중, 충청/호남 등 지방 시장 침투 확대 요구됨")
    
    # 표 데이터
    regions = [
        ("수도권 (서울/경기/인천)", "15억 8,410만 원", "62.3%"),
        ("영남권 (부산/대구/경남/경북)", "4억 7,520만 원", "18.7%"),
        ("충청권 (대전/세종/충남/충북)", "2억 3,820만 원", "9.4%"),
        ("호남권 (광주/전남/전북)", "1억 5,740만 원", "6.2%"),
        ("기타 (강원/제주)", "8,710만 원", "3.4%"),
        ("합계 (보고서 명시 기준)", "25억 4,200만 원", "100.0%")
    ]
    
    # 좌측 영역: 깔끔한 데이터 테이블 표 삽입
    rows = 7
    cols = 3
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.5)
    
    table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 컬럼 너비 조정
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.1)
    table.columns[2].width = Inches(1.5)
    
    # 헤더 설정
    headers = ["지역명", "매출액", "매출 비중"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        format_para(p, size_pt=14, bold=True, color_rgb=COLOR_WHITE, align=PP_ALIGN.CENTER)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
    # 데이터 행 추가
    for r, (reg, val, pct) in enumerate(regions):
        data_row = [reg, val, pct]
        for c, val_str in enumerate(data_row):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            # 마지막 합계 행은 회색조로 하이라이트
            if r == len(regions) - 1:
                cell.fill.fore_color.rgb = RGBColor(230, 235, 240)
                font_bold = True
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE
                font_bold = False
                
            p = cell.text_frame.paragraphs[0]
            p.text = val_str
            # 지역명은 좌측정렬, 숫자는 우측정렬
            align_t = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            # 합계 행 글자 진하게
            format_para(p, size_pt=12, bold=font_bold, color_rgb=COLOR_DARK, align=align_t)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
    # 우측 영역: 분석 및 시사점 설명 카드
    create_card(slide6, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.5))
    
    t_box = slide6.shapes.add_textbox(Inches(8.2), Inches(2.0), Inches(4.1), Inches(4.1))
    tf = t_box.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = "지역별 현황 분석 및 시사점"
    format_para(p_title, size_pt=18, bold=True, color_rgb=COLOR_PRIMARY)
    p_title.space_after = Pt(15)
    
    insights = [
        "■ 수도권 편중 현상 심화 (+62%)",
        "  · 매출의 과반수가 서울·경기·인천에 집중됨.",
        "  · 대기업 및 IT 인프라 인프라 수혜 영향.",
        "■ 비수도권 잠재 시장 개척 필요",
        "  · 영남(18.7%) 대비 충청(9.4%), 호남(6.2%)의 매출 점유율이 심각하게 저조함.",
        "  · 호남·강원 지역 파트너망 협력 강화를 통해 침투율 개선 필요."
    ]
    
    for ins in insights:
        p_ins = tf.add_paragraph()
        p_ins.text = ins
        if ins.startswith("■"):
            format_para(p_ins, size_pt=13, bold=True, color_rgb=COLOR_DARK)
            p_ins.space_before = Pt(8)
        else:
            format_para(p_ins, size_pt=11, bold=False, color_rgb=COLOR_MUTED)
            p_ins.space_after = Pt(3)
            
    add_notes(slide6, 
        "이 슬라이드에서 말씀드릴 것은 상반기 지역별 매출 분포 현황과 주요 분석 내용입니다. "
        "구체적으로 말씀드리면, 당사 매출의 62.3%인 15억 8,410만 원이 서울, 경기, 인천 등 수도권에 집중되어 있습니다. "
        "이는 주요 기업 본사의 지리적 여건과 B2B 인프라가 수도권에 밀집되어 있는 데 기인합니다. "
        "반면 영남권은 18.7%로 뒤를 잇고 있지만 충청권 9.4%, 호남권 6.2%, 강원 및 제주는 3.4% 수준으로 지방 지사의 매출 비중이 매우 저조합니다. "
        "이처럼 수도권에 대한 매출 의존도가 극도로 높기 때문에, 당사가 하반기에 지속 성장하기 위해서는 호남과 강원 등 잠재력 대비 침투가 부진한 지방 시장의 파트너사 협력을 확대해야 합니다. "
        "다음 슬라이드에서는 이러한 실적을 바탕으로 도출한 주요 성과와 극복해야 할 이슈를 정리해 보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 7: 주요 성과 및 이슈 (Achievements & Issues)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "5. 주요 성과 및 이슈", "NPS 만족도 상승 및 대형 계약 수주 성공 vs 3월 비수기 급감 및 단가 하락 과제")
    
    # 2열 카드 배치 (좌측 성과 - 파란색 테두리/우측 이슈 - 오렌지 테두리)
    col_w = Inches(5.6)
    card_h = Inches(4.5)
    
    # 성과 카드
    create_card(slide7, Inches(0.8), Inches(1.8), col_w, card_h)
    # 타이틀바 (파랑)
    title_bar_perf = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), col_w, Inches(0.6))
    title_bar_perf.fill.solid()
    title_bar_perf.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    title_bar_perf.line.fill.background()
    p_tb_p = title_bar_perf.text_frame.paragraphs[0]
    p_tb_p.text = "✅ 주요 성과 (Key Achievements)"
    format_para(p_tb_p, size_pt=15, bold=True, color_rgb=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    title_bar_perf.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    perf_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.5), col_w - Inches(0.4), Inches(3.6))
    tf_perf = perf_box.text_frame
    tf_perf.word_wrap = True
    
    perfs = [
        "· 신규 고객 유치 목표 달성률 108% (목표 280건 → 실적 302건)",
        "· 고객 만족도 NPS 72점 기록 (전년 65점 대비 7점 상승)",
        "· 엔터프라이즈 대형 신규 계약 3건 수주 (연 5억+ 매출 기대)",
        "· 프리미엄 고객 업셀링 캠페인 ROI 340% 기록"
    ]
    for p in perfs:
        para = tf_perf.add_paragraph() if tf_perf.paragraphs[0].text else tf_perf.paragraphs[0]
        para.text = p
        format_para(para, size_pt=13, bold=False, color_rgb=COLOR_DARK)
        para.space_after = Pt(12)
        
    # 이슈 카드
    create_card(slide7, Inches(6.9), Inches(1.8), col_w, card_h)
    # 타이틀바 (오렌지)
    title_bar_iss = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.8), col_w, Inches(0.6))
    title_bar_iss.fill.solid()
    title_bar_iss.fill.fore_color.rgb = COLOR_LIGHT_ORANGE
    title_bar_iss.line.fill.background()
    p_tb_i = title_bar_iss.text_frame.paragraphs[0]
    p_tb_i.text = "⚠️ 극복 과제 (Key Issues)"
    format_para(p_tb_i, size_pt=15, bold=True, color_rgb=COLOR_SECONDARY, align=PP_ALIGN.CENTER)
    title_bar_iss.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    iss_box = slide7.shapes.add_textbox(Inches(7.1), Inches(2.5), col_w - Inches(0.4), Inches(3.6))
    tf_iss = iss_box.text_frame
    tf_iss.word_wrap = True
    
    issues = [
        "· 3월 분기 전환기 영업 공백으로 인한 비수기 매출 급감 발생",
        "· 베이직 요금제 고객 단가 하락 추세 (서비스 팩 도입 필요)",
        "· 호남, 강원 지역 침투율 저조로 전국 영업망 확장성 정체",
        "· 일부 엔터프라이즈 리드 타임 지연 대응 체계 구축 요망"
    ]
    for i in issues:
        para = tf_iss.add_paragraph() if tf_iss.paragraphs[0].text else tf_iss.paragraphs[0]
        para.text = i
        format_para(para, size_pt=13, bold=False, color_rgb=COLOR_DARK)
        para.space_after = Pt(12)
        
    add_notes(slide7, 
        "이 슬라이드에서 말씀드릴 것은 상반기 주요 성과 및 우리가 안고 있는 극복 과제입니다. "
        "구체적으로 말씀드리면, 성과 부문에서는 마케팅과 영업의 호조로 신규 유치 목표 대비 108%인 302건을 달성했고, "
        "지속적인 품질 개선과 적극적 지원을 통해 고객 만족 지수 NPS가 전년 대비 7점이나 향상된 72점을 마크했습니다. "
        "또한 향후 연간 5억 원 규모의 계약 가치가 예상되는 엔터프라이즈 3개사 신규 수주와 업셀링 캠페인의 높은 ROI 역시 고무적입니다. "
        "하지만 극복해야 할 이슈 또한 뚜렷합니다. 1분기 전환기인 3월 비수기의 영업 공백 관리 체계 부재로 매출이 급감했으며, "
        "베이직 요금제 고객의 평균 단가가 하락세에 있어 번들 부가 서비스 발굴이 시급합니다. "
        "아울러 앞서 보신 것처럼 호남과 강원 지역 침투율이 여전히 저조해 전국구 확장성에 애로가 있습니다. "
        "이러한 분석을 토대로 마련한 하반기 핵심 액션 제안을 보고드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 8: 하반기 전략 제안 (Future Strategies)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "6. 하반기 전략 제안", "상반기 성장 모멘텀 확대와 취약점 개선을 위한 4대 핵심 영업/마케팅 전략")
    
    # 2x2 그리드 레이아웃의 전략 카드 배치
    strategies = [
        ("01", "프리미엄 프로모션 강화", "상반기 최고 성장세를 보인 프리미엄 모멘텀 연장을 목표로 합니다.", 
         ["· 7~8월 여름 패키지 할인 (조기 연장 시 10% 추가 할인)", "· 9월 추석 맞춤형 기업 단체 선물 패키지 구성"]),
        ("02", "엔터프라이즈 관리 고도화", "대형 고객 유지를 통한 영업 수익 기반 확대를 공고히 합니다.", 
         ["· 기존 고객 이탈 방지를 위한 전용 세미나 분기별 운영", "· 하반기 중 대형 고객사 2개사 추가 수주 파이프라인 가동"]),
        ("03", "지역 파트너 협력 균형 전략", "매출 편중을 해결하고 지방 잠재 고객을 적극 유치합니다.", 
         ["· 호남 및 강원 지역 거점 파트너사 3곳 이상 계약 확보", "· 지방 공공기관 및 중소기업 디지털 전환 입찰 참여 확대"]),
        ("04", "디지털 마케팅 활성화", "온라인 잠재 고객 리드 확보를 위해 인프라를 증진합니다.", 
         ["· 마케팅 예산 20% 증액 편성으로 타겟 광고 고도화", "· 솔루션 온라인 실시간 데모 세션 주 1회 상설 정례화"])
    ]
    
    col_w = Inches(5.6)
    row_h = Inches(2.1)
    start_x = Inches(0.8)
    start_y = Inches(1.8)
    
    for i, (num, title, sub, bullets) in enumerate(strategies):
        r = i // 2
        c = i % 2
        lx = start_x + c * (col_w + Inches(0.5))
        ty = start_y + r * (row_h + Inches(0.3))
        
        # 전략 카드 생성 (중요 액션은 오렌지 톤 활용)
        bg_col = COLOR_LIGHT_ORANGE if i == 0 else COLOR_WHITE
        create_card(slide8, lx, ty, col_w, row_h, bg_color=bg_col)
        
        # 번호 데코
        num_box = slide8.shapes.add_textbox(lx + Inches(0.15), ty + Inches(0.15), Inches(0.8), Inches(0.4))
        p_n = num_box.text_frame.paragraphs[0]
        p_n.text = f"Action {num}"
        format_para(p_n, size_pt=13, bold=True, color_rgb=COLOR_PRIMARY)
        
        # 타이틀
        title_box = slide8.shapes.add_textbox(lx + Inches(1.3), ty + Inches(0.15), col_w - Inches(1.5), Inches(0.4))
        p_t = title_box.text_frame.paragraphs[0]
        p_t.text = title
        format_para(p_t, size_pt=15, bold=True, color_rgb=COLOR_DARK)
        
        # 내용 블록
        content_box = slide8.shapes.add_textbox(lx + Inches(0.2), ty + Inches(0.6), col_w - Inches(0.4), Inches(1.4))
        tf_c = content_box.text_frame
        tf_c.word_wrap = True
        
        p_s = tf_c.paragraphs[0]
        p_s.text = sub
        format_para(p_s, size_pt=11, bold=True, color_rgb=COLOR_MUTED)
        p_s.space_after = Pt(6)
        
        for bullet in bullets:
            p_b = tf_c.add_paragraph()
            p_b.text = bullet
            format_para(p_b, size_pt=11, bold=False, color_rgb=COLOR_DARK)
            p_b.space_after = Pt(2)
            
    add_notes(slide8, 
        "이 슬라이드에서 말씀드릴 것은 하반기 매출 극대화와 지역 확장을 위한 네 가지 핵심 전략 제안입니다. "
        "구체적으로 말씀드리면, 첫째, 프리미엄의 성장 모멘텀을 연장하기 위해 7-8월 휴가철 패키지와 조기 갱신 10% 할인을 유도하고 9월 추석 특수 기업 패키지를 기획하겠습니다. "
        "둘째, 엔터프라이즈의 안정성을 높이기 위해 기존 대형 고객사 갱신율 95% 달성을 목표로 분기별 전용 세미나를 정례화하고 신규 파이프라인 2건 이상을 연내 완결 짓겠습니다. "
        "셋째, 침투가 저조한 호남과 강원 지역에 거점 파트너사 3곳 이상을 확보하여 지방 중소기업 및 공공 조달 입찰 참여 기회를 확대하겠습니다. "
        "넷째, 효과적인 리드 발굴을 위해 디지털 마케팅 예산을 20% 증액하여 주 1회 온라인 라이브 데모 세션을 대대적으로 가동하겠습니다. "
        "이러한 전략들이 완벽히 수행된다면 하반기 목표 실적 또한 초과 달성할 것으로 보입니다. "
        "그럼 발표를 마무리하도록 하겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 9: 마무리 및 Q&A (Closing)
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(slide_layout)
    
    # 다크 배경으로 끝인사 분위기 환기
    bg_shape9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_shape9.fill.solid()
    bg_shape9.fill.fore_color.rgb = COLOR_PRIMARY
    bg_shape9.line.fill.background()
    
    # 데코 포인트 (오렌지 띠)
    deco_shape9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))
    deco_shape9.fill.solid()
    deco_shape9.fill.fore_color.rgb = COLOR_SECONDARY
    deco_shape9.line.fill.background()
    
    # 중앙 텍스트박스
    center_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf9 = center_box.text_frame
    tf9.word_wrap = True
    
    p9_title = tf9.paragraphs[0]
    p9_title.text = "경청해 주셔서 감사합니다."
    format_para(p9_title, size_pt=42, bold=True, color_rgb=COLOR_WHITE, align=PP_ALIGN.CENTER)
    p9_title.space_after = Pt(24)
    
    p9_sub = tf9.add_paragraph()
    p9_sub.text = "Q&A 및 자유 토론"
    format_para(p9_sub, size_pt=20, bold=True, color_rgb=COLOR_SECONDARY, align=PP_ALIGN.CENTER)
    p9_sub.space_after = Pt(24)
    
    p9_info = tf9.add_paragraph()
    p9_info.text = "문의: 영업기획팀 김민서 대리 (minseo.kim@company.com / 02-1234-5678)"
    format_para(p9_info, size_pt=14, bold=False, color_rgb=RGBColor(210, 225, 250), align=PP_ALIGN.CENTER)
    
    add_notes(slide9, 
        "이상으로 2025년 상반기 매출 보고 요약 발표를 마치겠습니다. "
        "오늘 공유드린 성과와 추진 과제, 그리고 하반기 네 가지 실행안에 대해 "
        "의문사항이나 추가 조언이 있으시다면 편안하게 말씀해 주시기 바랍니다. "
        "성심껏 답변해 드리도록 하겠습니다. 경청해 주셔서 감사합니다.")

    # 저장 경로
    output_path = "c:/Users/user/Desktop/ai사무자동화/예제파일/12_발표_프롬프트_실습/02_발표_재료/상반기매출_발표.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
