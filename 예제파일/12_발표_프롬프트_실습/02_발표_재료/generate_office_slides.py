import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR

# UTF-8 출력 재설정 (윈도우 한글 인코딩 깨짐 방지)
sys.stdout.reconfigure(encoding='utf-8')

# 색상 상수 정의
COLOR_BG = RGBColor(255, 255, 255)         # 배경 흰색
COLOR_PRIMARY = RGBColor(43, 87, 151)       # 파랑 (#2B5797)
COLOR_SECONDARY = RGBColor(232, 119, 34)    # 포인트 주황 (#E87722)
COLOR_DARK = RGBColor(43, 87, 151)          # 메인 글씨 색상 (파랑 통일)
COLOR_MUTED = RGBColor(120, 130, 145)       # 보조 그레이 색상 (푸터/쪽번호용)
COLOR_CARD_BG = RGBColor(245, 247, 250)     # 카드 배경용 연한 회색 (Off-white)
COLOR_WHITE = RGBColor(255, 255, 255)

def clean_slide(slide):
    """슬라이드의 모든 기존 요소를 삭제하고 배경을 흰색으로 초기화"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG
    
    shapes = list(slide.shapes)
    for s in shapes:
        sp = s._element
        sp.getparent().remove(sp)

def add_header(slide, tag_text, title_text, subtitle_text):
    """차별화된 헤더 디자인: 좌측 상단 카테고리 태그 블록 + 우측 메인 타이틀"""
    # 1. 좌측 카테고리 태그 도형 (파랑 또는 주황 채우기)
    tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(1.8), Inches(0.45))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = COLOR_PRIMARY
    tag_box.line.fill.background()
    
    tf_tag = tag_box.text_frame
    tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag_text
    p_tag.font.name = '맑은 고딕'
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_WHITE
    p_tag.alignment = PP_ALIGN.CENTER
    
    # 2. 타이틀 텍스트박스 (태그 옆에 나란히 배치)
    title_box = slide.shapes.add_textbox(Inches(2.8), Inches(0.35), Inches(9.7), Inches(0.55))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = '맑은 고딕'
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 3. 서브 메시지 텍스트박스 (주황색 포인트 활용)
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.92), Inches(11.7), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_top = tf_sub.margin_bottom = tf_sub.margin_left = tf_sub.margin_right = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.name = '맑은 고딕'
    p_sub.font.size = Pt(14)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_SECONDARY
    
    # 4. 하단 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 224, 230)
    line.line.fill.background()

def add_footer_and_pagenum(slide, page_num):
    """슬라이드 하단에 부서명("사무자동화 추진팀") 및 쪽번호 추가 (2cm 여백 고려: 6.9 inches)"""
    # 부서명 (좌측)
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(4.0), Inches(0.3))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "사무자동화 추진팀"
    p_l.font.name = '맑은 고딕'
    p_l.font.size = Pt(11)
    p_l.font.color.rgb = COLOR_MUTED
    
    # 쪽번호 (우측)
    right_box = slide.shapes.add_textbox(Inches(11.0), Inches(6.9), Inches(1.5), Inches(0.3))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = str(page_num)
    p_r.font.name = '맑은 고딕'
    p_r.font.size = Pt(11)
    p_r.font.color.rgb = COLOR_MUTED
    p_r.alignment = PP_ALIGN.RIGHT

def format_text_box(slide, text_lines, left, top, width, height, font_size=18, bg_color=None, border_color=None):
    """3줄 이하 텍스트 본문 생성 상자 (테두리가 있는 얇은 화이트 카드 지원)"""
    card = None
    if bg_color or border_color:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color if bg_color else COLOR_WHITE
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
            
    t_box = slide.shapes.add_textbox(left + Inches(0.2) if bg_color or border_color else left, 
                                     top + Inches(0.2) if bg_color or border_color else top, 
                                     width - Inches(0.4) if bg_color or border_color else width, 
                                     height - Inches(0.3) if bg_color or border_color else height)
    tf = t_box.text_frame
    tf.word_wrap = True
    
    for idx, line in enumerate(text_lines):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = line
        p.font.name = '맑은 고딕'
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_DARK
        p.font.bold = False
        p.space_after = Pt(12)

def add_notes(slide, notes_text):
    """슬라이드 노트 등록"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def main():
    prs = Presentation()
    # 16:9 와이드스크린 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # SLIDE 1: 표지 (디자인 변경: 미니멀 좌측 정렬 레이아웃 및 굵은 파랑 바 장식)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(slide_layout)
    clean_slide(slide1)
    
    # 좌측 장식 바
    left_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.18), Inches(4.0))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLOR_PRIMARY
    left_bar.line.fill.background()
    
    # 제목
    title_box = slide1.shapes.add_textbox(Inches(1.3), Inches(2.0), Inches(11.0), Inches(1.5))
    tf1 = title_box.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "사무 자동화 추진 사업 보고"
    p1.font.name = '맑은 고딕'
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    
    # 부제
    sub_box = slide1.shapes.add_textbox(Inches(1.3), Inches(3.3), Inches(11.0), Inches(1.0))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "생산성 혁신을 위한 단순 반복 업무 자동화 및 기대 효과"
    p_sub.font.name = '맑은 고딕'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_SECONDARY
    
    # 작성자 정보
    info_box = slide1.shapes.add_textbox(Inches(1.3), Inches(4.8), Inches(11.0), Inches(1.0))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "작성부서: 사무자동화 추진팀  |  작성일: 2026년 6월 22일"
    p_info.font.name = '맑은 고딕'
    p_info.font.size = Pt(14)
    p_info.font.color.rgb = COLOR_MUTED
    
    add_notes(slide1, 
        "안녕하십니까. 사무자동화 추진팀 발표를 맡은 전략기획팀 김민서 대리입니다. "
        "지금부터 사무 자동화 추진 사업 보고를 시작하겠습니다. "
        "이번 보고는 단순 반복 작업의 비중을 낮춰 실질적인 업무 효율과 생산성을 높이기 위한 전략적 방안을 논의하고자 준비했습니다. "
        "다음 슬라이드에서 먼저 추진 배경과 현황 분석을 살펴보겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 2: 추진 배경 및 현황 분석 (본문 1 - 3줄 이내)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    clean_slide(slide2)
    add_header(slide2, "추진 배경", "1. 추진 배경 및 현황 분석", "단순 반복 업무 비중 35% 도달로 인한 비효율 개선 시급")
    add_footer_and_pagenum(slide2, 2)
    
    backgrounds = [
        "· 비즈니스 환경: 최근 1년간 온라인 사무 자동화 도구 도입이 다양한 업계에서 급속하게 확산됨",
        "· 현황 분석 결과: 단순 반복성 행정/정리 업무가 전체 사무 처리 비중의 35%를 차지하여 비효율 초래",
        "· 핵심 개선 방향: 임직원의 문서 작성, 자료 정리 및 발표 보고서 제작 소요 시간을 낮추는 것이 핵심"
    ]
    format_text_box(slide2, backgrounds, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), font_size=18, bg_color=COLOR_CARD_BG)
    
    add_notes(slide2, 
        "이 슬라이드에서 말씀드릴 것은 사무 자동화 추진의 배경과 현황 분석입니다. "
        "구체적으로 말씀드리면, 최근 1년간 온라인 사무 자동화 도구가 빠르게 확산되면서 생산성을 높이려는 시장의 요구가 커졌습니다. "
        "자체 분석 결과, 전체 업무 중 단순 반복 작업이 약 35%를 차지하고 있으며, 이로 인해 시간적 영업 공백이 지속적으로 발생하는 것을 확인했습니다. "
        "따라서 문서 작성과 자료 정리에 드는 시간을 획기적으로 낮추는 것을 핵심 과제로 설정했습니다. "
        "다음 슬라이드에서는 이를 개선하기 위한 구체적인 세부 추진 방안을 설명해 드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 3: 사무 자동화 개선 방안 (본문 2 - 3단 화이트 카드 레이아웃)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    clean_slide(slide3)
    add_header(slide3, "개선 방안", "2. 사무 자동화 개선 방안", "초안 자동 생성, 파일 정리 자동화 및 템플릿 표준화 도입")
    add_footer_and_pagenum(slide3, 3)
    
    plans = [
        ("① AI 문서 초안 자동 생성", ["· 생성형 AI 기반 문서 초안 신속 작성", "· 핵심 키워드 중심 자동 텍스트 렌더링", "· 문서 개발 소요 시간 평균 60% 이상 감축"]),
        ("② 파일 정리·변환 자동화", ["· 확장자(csv, pdf) 일괄 자동 재패키징", "· 부서별 수작업 수집 경로의 파이썬 이관", "· 폴더 분류 및 파일 분류 수작업 제거"]),
        ("③ 템플릿 표준화", ["· 사내 문서 양식 및 발표 템플릿 통일", "· 1장 1메시지 슬라이드 작성 규격화", "· 가이드라인을 통한 불필요 서식 설계 제거"])
    ]
    
    col_w = Inches(3.7)
    card_h = Inches(4.5)
    start_x = Inches(0.8)
    spacing = Inches(0.3)
    
    for idx, (title, bullets) in enumerate(plans):
        lx = start_x + idx * (col_w + spacing)
        # 파란색 테두리가 있는 깔끔한 흰색 카드 레이아웃 생성
        format_text_box(slide3, [title] + bullets, lx, Inches(1.8), col_w, card_h, font_size=13, bg_color=COLOR_WHITE, border_color=COLOR_PRIMARY)
        
    add_notes(slide3, 
        "이 슬라이드에서 말씀드릴 것은 구체적인 사무 자동화 3대 개선 방안입니다. "
        "구체적으로 말씀드리면, 첫째, AI를 연동한 문서 초안 작성을 전면 도입해 초고 생성 시간을 단축하겠습니다. "
        "둘째, 사람이 일일이 폴더를 뒤적이며 정리하던 작업과 확장자 변환 과정을 파이썬 스크립트 기반으로 자동화하겠습니다. "
        "셋째, 사내 보고서와 발표 장표 템플릿 표준화를 선언해 쓸데없는 서식 꾸미기 시간을 없애겠습니다. "
        "다음 슬라이드에서는 이 방안들을 적용했을 때 얻을 수 있는 구체적인 기대 효과 수치를 막대 그래프로 보여드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 4: 도입 기대 효과 및 시간 절감 (본문 3 - 세로 막대 차트 삽입)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    clean_slide(slide4)
    add_header(slide4, "기대 효과", "3. 도입 기대 효과 및 시간 절감", "사무 자동화 도입 전/후 대비 연간 약 1,200시간 업무 효율 향상")
    add_footer_and_pagenum(slide4, 4)
    
    # 차트 데이터 준비 (도입 전/후 시간 비교)
    chart_data_col = CategoryChartData()
    chart_data_col.categories = ['문서 작성 (시간)', '자료 정리 (시간)', '보고서 제작 (시간)']
    chart_data_col.add_series('도입 전', (8.0, 5.0, 6.0))
    chart_data_col.add_series('도입 후', (3.0, 1.0, 2.0))
    
    # 세로 막대 차트 삽입 (Clustered Column)
    x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.5)
    chart_col = slide4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_col
    ).chart
    chart_col.has_legend = True
    chart_col.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # 우측 텍스트 설명 카드 (테두리 적용)
    format_text_box(slide4, [
        "■ 도입 전/후 대비 절감률",
        "· 문서 작성 시간: 8h → 3h (62% 절감)",
        "· 자료 정리 시간: 5h → 1h (80% 절감)",
        "· 보고서 제작 시간: 6h → 2h (67% 절감)",
        "※ 연간 약 1,200시간 사무 비용 축소"
    ], Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.5), font_size=13, bg_color=COLOR_WHITE, border_color=COLOR_SECONDARY)
    
    add_notes(slide4, 
        "이 슬라이드에서 말씀드릴 것은 자동화 3대 방안의 도입 전후 대비 업무 시간 절감 수치입니다. "
        "구체적으로 말씀드리면, 세로 막대 차트에서 볼 수 있듯이 기존에 평균 8시간 소요되던 문서 작성이 3시간으로 줄어 62%가 절감됩니다. "
        "또한 5시간 걸리던 자료 정리는 1시간으로 단축되어 80%의 높은 효율을 냅니다. "
        "마지막으로 6시간 걸리던 보고서 제작도 2시간으로 개선되어 67% 절감을 이루며, 결과적으로 연간 1,200시간 가량의 자원을 아끼는 성과를 거둘 것으로 보입니다. "
        "그럼 마지막 장표로 넘어가 핵심 요약 및 마무리 안내를 해 드리겠습니다.")

    # -------------------------------------------------------------
    # SLIDE 5: 요약 및 마무리 (요약/제안 및 Q&A 통합 1장)
    # 가이드: 발견 3줄 / 제안 2줄 + Q&A 안내
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    clean_slide(slide5)
    add_header(slide5, "요약/마무리", "핵심 요약 및 차기 프로세스 실행 계획 제안", "연간 1,200시간 생산성 증진 및 프로세스 표준 가동")
    add_footer_and_pagenum(slide5, 5)
    
    # 좌측: 핵심 발견 3줄 (주황 테두리 카드)
    format_text_box(slide5, [
        "[핵심 발견 및 효과]",
        "1. 단순 행정 반복 업무 35%에 대한 자동화 적용으로 부서별 효율성 확보",
        "2. AI 문서 초안 작성 도입을 통해 문서 작성 소요 시간 최대 62% 단축",
        "3. 파일 자동 변환 및 양식 규격화 도입으로 연간 총 1,200시간 리소스 확보"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5), font_size=12, bg_color=COLOR_WHITE, border_color=COLOR_SECONDARY)
    
    # 우측: 제안 2줄 및 Q&A 안내 (파랑 테두리 카드)
    format_text_box(slide5, [
        "[제안 및 향후 조치]",
        "1. 7월 2째주까지 부서별 자동화 템플릿 표준 보급 및 파일 변환 가이드라인 배포",
        "2. 추진 단계별 현업 피드백 수집을 위한 분기별 사용자 모니터링 체계 가동",
        "",
        "※ Q&A 및 토론 문의처:",
        "사무자동화 추진팀 (oa_team@company.com / 내선 200)"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5), font_size=12, bg_color=COLOR_WHITE, border_color=COLOR_PRIMARY)
    
    add_notes(slide5, 
        "이 슬라이드에서 말씀드릴 것은 하반기 비즈니스의 종합 제언입니다. "
        "구체적으로 말씀드리면, 수작업 비중 35%를 개선해 연간 1,200시간의 여유 리소스를 확보하는 것을 기대하고 있으며, "
        "이를 수행하기 위해 부서별 표준 보급 및 피드백 모니터링을 즉시 단행할 계획입니다. "
        "이상으로 발표를 마치고 질의응답 세션을 시작하도록 하겠습니다. 질문이나 좋은 의견이 있으시면 편하게 개진해 주시기 바랍니다. "
        "감사합니다.")
    
    # 저장
    output_path = "c:/Users/user/Desktop/ai사무자동화/예제파일/12_발표_프롬프트_실습/02_발표_재료/사무자동화_추진_보고.pptx"
    prs.save(output_path)
    print(f"Office automation presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
